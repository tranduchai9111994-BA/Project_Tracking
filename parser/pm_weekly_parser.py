# -*- coding: utf-8 -*-
"""
Parser Weekly Report PPT (python-pptx).

Trích text + bảng từ slides — không pixel-perfect:
  - Cover: tiêu đề + khoảng ngày báo cáo
  - Done this week: bảng STT / Công việc / Đơn vị / Ngày / Tình trạng / Ghi chú
  - Next week: bảng STT / Công việc / Đơn vị / Ngày bắt đầu / Ngày kết thúc
  - Issues / Risks: text hoặc N/A
"""
from __future__ import annotations

import re
from datetime import datetime
from typing import Any, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _norm(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").strip())


def _parse_vn_date(s: str) -> Optional[str]:
    s = (s or "").strip()
    if not s:
        return None
    for fmt in ("%d/%m/%Y", "%d-%m-%Y", "%Y-%m-%d", "%d.%m.%Y"):
        try:
            return datetime.strptime(s[:10], fmt).date().isoformat()
        except ValueError:
            continue
    return None


def _extract_period(texts: list[str]) -> tuple[Optional[str], Optional[str]]:
    """Tìm khoảng 'dd/mm/yyyy - dd/mm/yyyy' hoặc 'từ ngày A đến B'."""
    blob = "\n".join(texts)
    m = re.search(
        r"(\d{1,2}[/.-]\d{1,2}[/.-]\d{2,4})\s*[-–—]\s*(\d{1,2}[/.-]\d{1,2}[/.-]\d{2,4})",
        blob,
    )
    if m:
        return _parse_vn_date(m.group(1)), _parse_vn_date(m.group(2))
    m2 = re.search(
        r"từ\s+ngày\s+(\d{1,2}[/.-]\d{1,2}[/.-]\d{2,4})\s+đến\s+(\d{1,2}[/.-]\d{1,2}[/.-]\d{2,4})",
        blob,
        re.IGNORECASE,
    )
    if m2:
        return _parse_vn_date(m2.group(1)), _parse_vn_date(m2.group(2))
    return None, None


def _shape_texts(shape) -> list[str]:
    if not getattr(shape, "has_text_frame", False):
        return []
    out = []
    for para in shape.text_frame.paragraphs:
        t = para.text.strip()
        if t:
            out.append(t)
    return out


def _table_rows(shape) -> list[list[str]]:
    table = shape.table
    rows = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def _header_kind(header: list[str]) -> str:
    h = " | ".join(header).lower()
    if "tuần tiếp" in h or "tuan tiep" in h or ("bắt đầu" in h and "kết thúc" in h):
        return "next"
    if "tình trạng" in h or "tinh trang" in h or "công việc trong tuần" in h:
        return "done"
    if "rủi ro" in h or "rui ro" in h or "vấn đề" in h or "issue" in h:
        return "risk"
    return "other"


def _rows_to_done(rows: list[list[str]]) -> list[dict]:
    if not rows:
        return []
    header = rows[0]
    # Map by header keywords
    idx = {i: _norm(h).lower() for i, h in enumerate(header)}
    def find(*keys):
        for i, h in idx.items():
            if any(k in h for k in keys):
                return i
        return None
    i_stt = find("stt") or 0
    i_task = find("công việc", "cong viec") or 1
    i_unit = find("đơn vị", "don vi") or 2
    i_date = find("ngày", "ngay") or 3
    i_status = find("tình trạng", "tinh trang", "status") or 4
    i_note = find("ghi chú", "ghi chu", "note")
    items = []
    for row in rows[1:]:
        if not any(cell.strip() for cell in row):
            continue
        task = row[i_task] if i_task < len(row) else ""
        if not task:
            continue
        items.append({
            "stt": row[i_stt] if i_stt < len(row) else "",
            "task": task,
            "unit": row[i_unit] if i_unit < len(row) else "",
            "date": row[i_date] if i_date < len(row) else "",
            "date_iso": _parse_vn_date(row[i_date] if i_date < len(row) else ""),
            "status": row[i_status] if i_status < len(row) else "",
            "note": row[i_note] if i_note is not None and i_note < len(row) else "",
        })
    return items


def _rows_to_next(rows: list[list[str]]) -> list[dict]:
    if not rows:
        return []
    header = rows[0]
    idx = {i: _norm(h).lower() for i, h in enumerate(header)}
    def find(*keys):
        for i, h in idx.items():
            if any(k in h for k in keys):
                return i
        return None
    i_stt = find("stt") or 0
    i_task = find("công việc", "cong viec") or 1
    i_unit = find("đơn vị", "don vi") or 2
    i_start = find("bắt đầu", "bat dau", "start") or 3
    i_end = find("kết thúc", "ket thuc", "end") or 4
    i_note = find("ghi chú", "ghi chu", "note")
    items = []
    for row in rows[1:]:
        if not any(cell.strip() for cell in row):
            continue
        task = row[i_task] if i_task < len(row) else ""
        if not task:
            continue
        start = row[i_start] if i_start < len(row) else ""
        end = row[i_end] if i_end < len(row) else ""
        items.append({
            "stt": row[i_stt] if i_stt < len(row) else "",
            "task": task,
            "unit": row[i_unit] if i_unit < len(row) else "",
            "start": start,
            "end": end,
            "start_iso": _parse_vn_date(start),
            "end_iso": _parse_vn_date(end),
            "note": row[i_note] if i_note is not None and i_note < len(row) else "",
        })
    return items


def preview_weekly(filepath: str) -> dict[str, Any]:
    """Tóm tắt slides + đề xuất section (không lưu)."""
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        tables = 0
        for shape in slide.shapes:
            if shape.has_table:
                tables += 1
            texts.extend(_shape_texts(shape))
        title = texts[0] if texts else f"Slide {i}"
        kind = _classify_slide(texts, tables)
        slides.append({
            "index": i,
            "title": title[:120],
            "kind": kind,
            "text_count": len(texts),
            "table_count": tables,
        })
    return {
        "slide_count": len(prs.slides),
        "slides": slides,
        "proposed_sections": {
            "done": [s["index"] for s in slides if s["kind"] == "done"],
            "next": [s["index"] for s in slides if s["kind"] == "next"],
            "risk": [s["index"] for s in slides if s["kind"] == "risk"],
            "issues": [s["index"] for s in slides if s["kind"] == "issues"],
            "cover": [s["index"] for s in slides if s["kind"] == "cover"],
        },
    }


def _classify_slide(texts: list[str], table_count: int) -> str:
    blob = "\n".join(texts).lower()
    if any("báo cáo tiến độ" in t.lower() or "bao cao tien do" in t.lower() for t in texts[:3]):
        return "cover"
    if "nội dung" in blob or "noi dung" in blob:
        return "toc"
    if re.search(r"^0\d$", texts[0].strip() if texts else ""):
        return "section"
    if "tiếp theo" in blob or "tiep theo" in blob:
        return "next" if table_count else "section"
    if "rủi ro" in blob or "rui ro" in blob or "vấn đề" in blob:
        return "risk"
    if "issues" in blob or "còn xử lý" in blob:
        return "issues"
    if "công việc trong tuần" in blob or "đã thực hiện" in blob:
        return "done"
    if table_count and ("tình trạng" in blob or "hoàn thành" in blob):
        return "done"
    if "thank" in blob:
        return "thanks"
    return "other"


def parse_weekly(filepath: str) -> dict[str, Any]:
    """Parse PPT weekly → dict chuẩn hoá."""
    prs = Presentation(filepath)
    all_texts: list[str] = []
    slides_summary: list[dict] = []
    done: list[dict] = []
    next_items: list[dict] = []
    issues: list[str] = []
    risks: list[str] = []
    period_start = period_end = None
    title = ""
    project_title = ""
    seen_done_keys: set[str] = set()

    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
            if shape.has_table:
                tables.append(_table_rows(shape))
            texts.extend(_shape_texts(shape))
        all_texts.extend(texts)
        kind = _classify_slide(texts, len(tables))
        slides_summary.append({
            "index": i,
            "title": (texts[0] if texts else f"Slide {i}")[:120],
            "kind": kind,
        })

        if i == 1 or kind == "cover":
            if not title and texts:
                title = texts[0]
            for t in texts:
                if re.search(r"\d{1,2}/\d{1,2}/\d{4}", t) and "-" in t:
                    ps, pe = _extract_period([t])
                    if ps:
                        period_start, period_end = ps, pe
                elif len(t) > 30 and ("dự án" in t.lower() or "ihrp" in t.lower()):
                    project_title = project_title or t

        # Period from body headings
        ps, pe = _extract_period(texts)
        if ps and not period_start:
            period_start, period_end = ps, pe

        for table in tables:
            if not table:
                continue
            hk = _header_kind(table[0])
            if hk == "done" or (kind == "done" and hk != "next"):
                for item in _rows_to_done(table):
                    key = f"{item['stt']}|{item['task'][:80]}"
                    if key not in seen_done_keys:
                        seen_done_keys.add(key)
                        done.append(item)
            elif hk == "next" or kind == "next":
                next_items.extend(_rows_to_next(table))

        # Issues / risks free text
        if kind in ("issues", "risk"):
            for t in texts:
                tl = t.lower().strip()
                if tl in ("n/a", "na", "-") or len(t) < 3:
                    if tl in ("n/a", "na"):
                        (issues if kind == "issues" else risks).append("N/A")
                    continue
                if any(x in tl for x in ("các vấn đề", "rủi ro", "issues", "công việc", "tiêu đề")):
                    continue
                if kind == "issues":
                    issues.append(t)
                else:
                    risks.append(t)

    # Dedup next by stt+task
    seen_next: set[str] = set()
    next_dedup = []
    for item in next_items:
        key = f"{item['stt']}|{item['task'][:80]}"
        if key not in seen_next:
            seen_next.add(key)
            next_dedup.append(item)

    return {
        "title": title or "Báo cáo tiến độ",
        "project_title": project_title,
        "period_start": period_start,
        "period_end": period_end,
        "done": done,
        "next": next_dedup,
        "issues": issues or ["N/A"],
        "risks": risks or ["N/A"],
        "slides_summary": slides_summary,
        "summary": {
            "slide_count": len(prs.slides),
            "done_count": len(done),
            "next_count": len(next_dedup),
            "has_issues": bool(issues) and issues != ["N/A"],
            "has_risks": bool(risks) and risks != ["N/A"],
        },
    }
