# -*- coding: utf-8 -*-
"""Tests chiều PM — parser + store + export (MVP)."""
from __future__ import annotations

import json
import os
from pathlib import Path

import openpyxl
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from analyzer import pm_store
from exporter.pm_exporter import export_pm_report
from parser.pm_plan_parser import parse_plan, propose_sheet_mapping, preview_plan_workbook
from parser.pm_weekly_parser import parse_weekly, preview_weekly


FIXTURES = Path(__file__).parent / "fixtures"
FIXTURES.mkdir(exist_ok=True)


def _make_plan_xlsx(path: Path) -> Path:
    wb = openpyxl.Workbook()
    # Gantt
    ws = wb.active
    ws.title = "Gantt-chart"
    ws["A1"] = "STT"
    ws["B1"] = "Công việc"
    ws["C1"] = "Lịch trình triển khai (Tuần)"
    ws["C2"] = "10/2025"
    ws["C3"] = "W1"
    ws["D3"] = "W2"
    ws["E3"] = "W3"
    ws["A4"] = 1
    ws["B4"] = "Khởi động dự án"
    ws["A5"] = 2
    ws["B5"] = "Khảo sát yêu cầu"
    # Schedule
    ws2 = wb.create_sheet("Lịch trình UAT_Golive")
    ws2["A1"] = "LỊCH TRÌNH"
    ws2["A2"] = "Công việc"
    ws2["B2"] = "Từ ngày"
    ws2["C2"] = "Đến ngày"
    ws2["D2"] = "Phụ trách chính FPT"
    ws2["E2"] = "Hỗ trợ FPT"
    ws2["F2"] = "Phụ trách chính MPHG"
    ws2["G2"] = "Hỗ trợ MPHG"
    ws2["H2"] = "Ghi chú"
    ws2["A3"] = "Giai đoạn 1: Golive demo"
    ws2["B3"] = "2025-12-22"
    ws2["C3"] = "2026-01-15"
    ws2["A4"] = "Quản trị hệ thống"
    ws2["B4"] = "2025-12-22"
    ws2["C4"] = "2025-12-26"
    ws2["D4"] = "NhiVN, ThiTM3"
    ws2["E4"] = "BinhTT33"
    ws2["F4"] = "chị Huyền"
    # Deliverables
    ws3 = wb.create_sheet("Sản phẩm bàn giao")
    ws3["A2"] = "STT"
    ws3["B2"] = "Sản phẩm bàn giao"
    ws3["C2"] = "Ngày bàn giao"
    ws3["D2"] = "Loại"
    ws3["A3"] = "Khởi động"
    ws3["A4"] = 1
    ws3["B4"] = "Tài liệu kickoff"
    ws3["C4"] = "2025-10-31"
    ws3["D4"] = "Bản mềm"
    # Teams
    ws4 = wb.create_sheet("Đội dự án FPT")
    ws4["A1"] = "STT"
    ws4["B1"] = "Họ và tên"
    ws4["C1"] = "Chức vụ"
    ws4["D1"] = "Trách nhiệm"
    ws4["E1"] = "Email"
    ws4["A2"] = "Quản lý"
    ws4["A3"] = 1
    ws4["B3"] = "Thái Tâm Bình"
    ws4["C3"] = "PM"
    ws4["E3"] = "BinhTT33@fpt.com"
    ws5 = wb.create_sheet("Đội dự án MPHG")
    ws5["A1"] = "STT"
    ws5["B1"] = "Họ và tên"
    ws5["C1"] = "Chức vụ"
    ws5["D1"] = "Vai trò"
    ws5["E1"] = "Trách nhiệm"
    ws5["F1"] = "Email"
    ws5["A2"] = "Ban quản lý"
    ws5["A3"] = 1
    ws5["B3"] = "Nguyễn Minh Tuấn"
    ws5["C3"] = "PM"
    ws5["F3"] = "minhtuan@minhphu.com"
    wb.save(path)
    return path


def _make_weekly_pptx(path: Path) -> Path:
    prs = Presentation()
    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "BÁO CÁO TIẾN ĐỘ"
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = "15/12/2025 - 19/12/2025\nDỰ ÁN MUA BẢN QUYỀN FPT.iHRP"
    # Done table slide
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide2 = prs.slides.add_slide(blank)
    # Add a table
    rows, cols = 3, 6
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)
    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    headers = ["STT", "Công việc trong tuần", "Đơn vị", "Ngày", "Tình trạng", "Ghi chú"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "Hoàn tất tài liệu URD"
    table.cell(1, 2).text = "FPT, MPHG"
    table.cell(1, 3).text = "18/12/2025"
    table.cell(1, 4).text = "Hoàn thành"
    table.cell(2, 0).text = "2"
    table.cell(2, 1).text = "Cấu hình hợp đồng"
    table.cell(2, 2).text = "FPT"
    table.cell(2, 3).text = "19/12/2025"
    table.cell(2, 4).text = "Hoàn thành"
    # Title text
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    tf.text = "2. CÁC CÔNG VIỆC TRONG TUẦN"
    p = tf.add_paragraph()
    p.text = "2.1 Công việc đã thực hiện từ ngày 15/12/2025 đến 19/12/2025"
    # Next week
    slide3 = prs.slides.add_slide(blank)
    tx3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
    tx3.text_frame.text = "4. CÁC CÔNG VIỆC TIẾP THEO"
    table2 = slide3.shapes.add_table(2, 6, Inches(0.5), Inches(1.2), Inches(9), Inches(1)).table
    next_h = ["STT", "Công việc tuần tiếp theo", "Đơn vị", "Ngày bắt đầu", "Ngày kết thúc", "Ghi chú"]
    for i, h in enumerate(next_h):
        table2.cell(0, i).text = h
    table2.cell(1, 0).text = "1"
    table2.cell(1, 1).text = "Đào tạo Key User"
    table2.cell(1, 2).text = "FPT, MPHG"
    table2.cell(1, 3).text = "24/12/2025"
    table2.cell(1, 4).text = "25/12/2025"
    # Risk slide
    slide4 = prs.slides.add_slide(blank)
    tx4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    tx4.text_frame.text = "3. CÁC VẤN ĐỀ, RỦI RO"
    p = tx4.text_frame.add_paragraph()
    p.text = "N/A"
    prs.save(str(path))
    return path


@pytest.fixture
def plan_file(tmp_path: Path) -> Path:
    return _make_plan_xlsx(tmp_path / "kehoach.xlsx")


@pytest.fixture
def weekly_file(tmp_path: Path) -> Path:
    return _make_weekly_pptx(tmp_path / "weekly.pptx")


def test_propose_sheet_mapping():
    names = [
        "Gantt-chart-old",
        "Lịch trình UAT_Golive",
        "Gantt-chart",
        "Sản phẩm bàn giao",
        "Đội dự án FPT",
        "Đội dự án MPHG",
    ]
    m = propose_sheet_mapping(names)
    assert m["Gantt-chart"] == "gantt"
    assert m["Gantt-chart-old"] == "gantt_old"
    assert m["Lịch trình UAT_Golive"] == "schedule"
    assert m["Sản phẩm bàn giao"] == "deliverables"
    assert m["Đội dự án FPT"] == "team_vendor"
    assert m["Đội dự án MPHG"] == "team_client"


def test_parse_plan(plan_file: Path):
    preview = preview_plan_workbook(str(plan_file))
    assert len(preview["sheets"]) == 5
    data = parse_plan(str(plan_file), preview["proposed_mapping"])
    assert data["summary"]["milestone_count"] >= 2
    assert data["summary"]["schedule_count"] >= 2
    assert any(s.get("is_phase_header") for s in data["schedule"])
    task = next(s for s in data["schedule"] if s["name"] == "Quản trị hệ thống")
    assert "NhiVN" in task["pic_fpt"]
    assert task["start"] == "2025-12-22"
    assert data["summary"]["team_vendor_count"] >= 1
    assert data["summary"]["deliverable_count"] >= 1


def test_parse_weekly(weekly_file: Path):
    preview = preview_weekly(str(weekly_file))
    assert preview["slide_count"] >= 3
    data = parse_weekly(str(weekly_file))
    assert data["summary"]["done_count"] >= 2
    assert data["summary"]["next_count"] >= 1
    assert data["period_start"] == "2025-12-15"
    assert data["done"][0]["status"] == "Hoàn thành"


def test_pm_store_and_export(tmp_path: Path, plan_file: Path, weekly_file: Path):
    project_dir = tmp_path / "proj"
    project_dir.mkdir()
    plan = parse_plan(str(plan_file))
    weekly = parse_weekly(str(weekly_file))
    pm_store.save_plan(str(project_dir), plan, source_filename="k.xlsx", source_path=str(plan_file))
    pm_store.save_weekly(str(project_dir), weekly, source_filename="w.pptx", source_path=str(weekly_file))
    bundle = pm_store.load_pm_bundle(str(project_dir))
    assert bundle["has_plan"] and bundle["has_weekly"]
    out = export_pm_report(
        bundle["plan"], bundle["weekly"],
        str(tmp_path / "out"),
        project_code="MPHG",
    )
    assert os.path.isfile(out)
    wb = openpyxl.load_workbook(out)
    assert "Lịch trình" in wb.sheetnames
    assert "Weekly Done" in wb.sheetnames


def test_pm_api_flow(flask_client, plan_file, weekly_file):
    """End-to-end qua Flask test client."""
    client = flask_client
    r = client.get("/api/projects")
    assert r.status_code == 200
    projects = r.get_json()
    if isinstance(projects, dict):
        slugs = [p["slug"] for p in projects.get("projects", projects.get("items", []))]
    else:
        slugs = [p["slug"] for p in projects]
    slug = slugs[0] if slugs else "default"

    with open(plan_file, "rb") as f:
        r = client.post(
            f"/api/projects/{slug}/pm/plan/preview",
            data={"file": (f, "kehoach.xlsx")},
            content_type="multipart/form-data",
        )
    assert r.status_code == 200, r.get_data(as_text=True)
    prev = r.get_json()
    assert prev["tmp_id"]
    r = client.post(
        f"/api/projects/{slug}/pm/plan/confirm",
        json={"tmp_id": prev["tmp_id"], "filename": "kehoach.xlsx", "sheet_mapping": prev["proposed_mapping"]},
    )
    assert r.status_code == 200, r.get_data(as_text=True)
    assert r.get_json()["ok"]

    with open(weekly_file, "rb") as f:
        r = client.post(
            f"/api/projects/{slug}/pm/weekly/preview",
            data={"file": (f, "weekly.pptx")},
            content_type="multipart/form-data",
        )
    assert r.status_code == 200, r.get_data(as_text=True)
    wprev = r.get_json()
    r = client.post(
        f"/api/projects/{slug}/pm/weekly/confirm",
        json={"tmp_id": wprev["tmp_id"], "filename": "weekly.pptx"},
    )
    assert r.status_code == 200, r.get_data(as_text=True)

    r = client.get(f"/api/projects/{slug}/pm")
    assert r.status_code == 200
    data = r.get_json()
    assert data["has_plan"] and data["has_weekly"]

    r = client.get(f"/api/projects/{slug}/pm/export")
    assert r.status_code == 200
    assert "spreadsheet" in (r.content_type or "") or r.data[:2] == b"PK"
